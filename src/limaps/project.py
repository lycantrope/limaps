import logging
import os
import sys
from dataclasses import dataclass, field
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any, Dict, Hashable, List, Optional, Tuple, Union

import coloredlogs
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import patches
from matplotlib.figure import Figure
from pptx import Presentation
from pptx.util import Cm, Pt

from .dotplot import add_jitter_plot
from .samplegroup import Samplegroup
from .utils import askdirectory

plt.rcParams["figure.max_open_warning"] = 100

UNDEFINE_LABEL = "unknown"
HEATSHOCK_LABELS = {"sis", "stress-induced sleep", "heatshock", "hs"}


@dataclass
class Project:

    foqthreshold: float = 0.2

    colnum: int = 8
    rownum: int = 6
    grouporder: int = "horizontal"
    mode: str = "lethargus"

    uniquegroupnames: List[str] = field(init=False, default_factory=list)
    groups: pd.MultiIndex = field(init=False, repr=False, default=None)
    manual_groups: List[str] = field(init=False, default_factory=list)

    date: str = field(init=False)
    expname: str = field(init=False)
    expnum: str = field(init=False)
    interval: float = field(init=False)

    datapath: Path = field(init=False, default=None)
    homepath: Path = field(init=False)

    data: pd.DataFrame = field(init=False, repr=False)

    samplegroups: List[Samplegroup] = field(
        init=False,
        repr=False,
        default_factory=list,
    )

    # internal value
    _summary: pd.DataFrame = field(init=False, repr=False, default=None)

    def __post_init__(self):
        self._analyzed_time = datetime.now().strftime(r"%y%m%d_%H%M%S")

    def __setstate__(self, kws: Dict[str, Any]) -> None:
        kws["datapath"] = Path(kws.get("datapath"))
        kws["homepath"] = Path(kws.get("homepath"))
        self.__dict__.update(kws)
        groups = kws["groups"]
        sgs = {sg.groupname: sg for sg in kws["samplegroups"]}
        datas = sorted(
            (
                (idx, sgs[name].fullindlist[ind_num - 1].rawdata)
                for name, ind_num, idx in groups
            ),
            key=lambda x: x[0],
        )
        df = pd.concat([d for _, d in datas], axis=1)
        df.columns = groups
        self.__dict__["data"] = df

    def __getstate__(self) -> Dict[str, Any]:
        kwargs = self.__dict__.copy()
        kwargs["datapath"] = os.fspath(kwargs["datapath"])
        kwargs["homepath"] = os.fspath(kwargs["homepath"])
        kwargs.pop("data")
        return kwargs

    def set_logger(self, level=logging.INFO) -> "Project":
        stream = logging.StreamHandler(sys.stdout)
        logging.basicConfig(
            format="%(asctime)s : %(module)s.%(funcName)s : %(levelname)s : %(message)s",
            level=level,
            handlers=[
                logging.FileHandler(self.datapath.parent.joinpath("history.log")),
                stream,
            ],
        )
        logger = logging.getLogger(__name__)
        coloredlogs.install(
            level=level,
            logger=logger,
            fmt="%(asctime)s.%(msecs)03d %(filename)s:%(lineno)d %(levelname)s %(message)s",
        )
        return self

    def set_targetfile(
        self,
        targetfile: Optional[Union[str, Path]] = None,
    ) -> "Project":

        targetfile = targetfile or askdirectory()
        if targetfile is None or not Path(targetfile).is_file():
            raise FileNotFoundError(f"File is not selected or not exist: {targetfile}")

        self.datapath = Path(targetfile)
        header = self.datapath.stem.split("_")
        if len(header) != 4 or not all(header):
            raise Exception(
                f"Format is wrong: {self.datapath.name}.\nPlease rename the filename as `date_groupname_expnum_interval.csv`"
            )

        foldername = f"{self._analyzed_time}_analyzed"
        if self.mode in HEATSHOCK_LABELS:
            foldername = f"{self._analyzed_time}_HS_analyzed"
        self.homepath = self.datapath.parent.joinpath(foldername)

        # setup logger
        self.set_logger()
        logging.info(
            f"""
Basic Parameters:
Threshold: {self.foqthreshold}
number of column: {self.colnum}
number of row: {self.rownum}
group orientation: {self.grouporder}
mode: {self.mode}
"""
        )

        self.date = header[0]
        self.expname = header[1]
        self.expnum = header[2]
        self.interval = float(header[3])
        logging.info(
            f"""
File information:
filename: {self.datapath.name}
date: {self.date}
experiment name: {self.expname}
number of experiment: {self.expnum}
interval between frame (sec): {self.interval}
"""
        )
        return self

    def set_groupnames(
        self, uniquegroupnames: List[str], gindex: List[Tuple[int, int]]
    ) -> "Project":
        if self.datapath is None:
            return self.set_targetfile().set_groupnames(uniquegroupnames, gindex)

        if len(gindex) != len(uniquegroupnames):
            raise Exception("The sizes of uniquegroupnames and gindex are not equal")
        if any(map(lambda x: not isinstance(x, (int, list, tuple)), gindex)):
            raise Exception(
                "The gindex must be List[int], List[Tuple[int]], List[Tuple[int,int]]"
            )

        self.uniquegroupnames = uniquegroupnames
        indexgrid = np.arange(self.colnum * self.rownum, dtype=int).reshape(
            self.rownum, self.colnum
        )
        sample_nums = self.rownum
        if self.grouporder.lower() in ("h", "horizontal"):
            indexgrid = np.fliplr(indexgrid.T)
            sample_nums = self.colnum
        groups = []
        remains = set(range(sample_nums))

        for idx, name in zip(gindex, uniquegroupnames):
            if isinstance(idx, int):
                ilist = [idx - 1]
            else:
                if len(idx) == 1:
                    ilist = [idx[0] - 1]
                else:
                    start, end = idx
                    ilist = list(range(start - 1, end))
            remains -= set(ilist)
            groups.extend(
                [
                    (name, i + 1, col)
                    for i, col in enumerate(indexgrid[ilist, :].ravel())
                ]
            )

        if remains:
            ilist = sorted(remains)
            groups.extend(
                [
                    (UNDEFINE_LABEL, i + 1, col)
                    for i, col in enumerate(indexgrid[ilist, :].ravel())
                ]
            )
            logging.warn(f"{ilist} are assigned as `unknown` group")

        group_grid = [["-" for _ in range(self.colnum)] for _ in range(self.rownum)]
        width = 0
        for name, i, num in groups:
            width = max(width, len(name))
            group_grid[num // self.colnum][num % self.colnum] = f"{name}_{i}"

        group_str = "\n" + "\n".join(
            ["".join(map(lambda s: s.ljust(width + 4), l)) for l in group_grid]
        )
        logging.info("\n Remi grid")
        logging.info(group_str)

        self.groups: Tuple[pd.MultiIndex, np.ndarray] = pd.MultiIndex.from_tuples(
            groups
        ).sortlevel(2)[0]
        return self

    def set_manual_groups(self, groupnames: List[str]) -> "Project":
        if not set(groupnames).issubset(self.uniquegroupnames):
            raise ValueError(f"{groupnames} has invalid value")
        self.manual_groups = groupnames

        foldername = f"{self._analyzed_time}_manually_analyzed"
        if self.mode in HEATSHOCK_LABELS:
            foldername = f"{self._analyzed_time}_HS_manually_analyzed"
        self.homepath = self.homepath.with_name(foldername)
        return self

    def read_dataframe(self) -> "Project":
        seps = {".csv": ",", ".xls": "\t", ".xlxs": "\t", ".txt": " "}
        sep = seps.get(self.datapath.suffix, ",")
        path = self.datapath
        df = pd.read_csv(path, sep=sep, nrows=2)
        # check auto_roi csv
        area_header = any(map(lambda col: "Area" in col, df.columns))

        if not area_header:
            df = pd.read_csv(path, sep=sep, header=None)
            df.columns = [f"Area.{col}" for col in df.columns]
        else:
            df = pd.read_csv(path, sep=sep)
        # extract area data only
        self.data = df[[col for col in df.columns if "area" in col.lower()]]
        self.data.columns = self.groups
        return self

    def init_a_samplegroup(self, df: pd.DataFrame) -> Samplegroup:
        groupname = df.columns.get_level_values(0)[0]

        sg = (
            Samplegroup(
                self.date,
                groupname,
                self.expnum,
                self.foqthreshold,
                is_dummy=groupname == UNDEFINE_LABEL,
            )
            .set_directory(self.homepath)
            .processindividuals(
                interval=self.interval,
                dataframe=df.sort_index(axis=1, level=1).droplevel(
                    level=[0, 1], axis=1
                ),
            )
        )
        if groupname in self.manual_groups:
            sg.manual_screen()
        return sg

    def batch_lethargus_samplegroup(self, df: pd.DataFrame):
        sg = self.init_a_samplegroup(df)

        sg = sg.saveafig(
            sg.makeallfigureofarea(
                60,
                sg.fullindlist,
            ),
            "rawareafull",
        ).saveafig(
            sg.plot_foq_heatmap(
                align=False,
            ),
            "foq_heatmap",
        )
        if sg.is_valid:
            (
                sg.makedf(ltlist=sg.fullindlist)
                .savesummarydf()
                .makeltmatrix(align="head")
                .saveltmatrix(sg.ltfoqmatrix, "fqlt")
                .makeltmatrix(align="tail")
                .saveltmatrix(sg.ltfoqmatrixtail, "fqltaligntail")
                .saveafig(sg.makealignfigure(sg.ltfoqmatrix, True, "mean"), "fqaligned")
                .saveafig(
                    sg.makealignfigure(sg.ltfoqmatrixtail, False, "mean"),
                    "fqalignedtail",
                )
                .saveafig(sg.makeallfigureofareacropped(sg.indlist), "areacropped")
                .saveafig(sg.plot_foq_heatmap(align=True), "foq_heatmap_aligned")
            )

        self.samplegroups.append(sg)
        return

    def batch_heatshock_samplegroup(self, df: pd.DataFrame):
        sg = self.init_a_samplegroup(df)

        self.samplegroups.append(
            sg.make_heatshock_df(start_hr=1.5, end_hr=3.5)
            .saveafig(
                sg.makeallfigureofarea(
                    60,
                    sg.fullindlist,
                ),
                "rawareafull",
            )
            .saveafig(
                sg.plot_foq_heatmap(
                    align=False,
                    xlim=timedelta(hours=5.0).seconds / self.interval,
                ),
                "foq_heatmap",
            )
        )

        return

    def process_samplegroups(self) -> "Project":
        if not self.homepath.exists():
            self.homepath.mkdir()

        if self.mode in HEATSHOCK_LABELS:
            self.data.groupby(axis=1, level=0).agg(self.batch_heatshock_samplegroup)
            return self

        self.data.groupby(axis=1, level=0).agg(self.batch_lethargus_samplegroup)
        return self

    def saveafig(
        self,
        figure: Figure,
        filename: str,
        *args: Tuple[Any],
        **kwargs: Dict[str, Any],
    ) -> "Project":
        figpath = self.homepath.joinpath(filename)
        figure.savefig(figpath, *args, **kwargs)
        return self

    def plot_samplegroups_grid(
        self,
        plotdata: str = "foq",
        window: int = 60,
        overlayparam: Optional[str] = None,
        meanduration: float = 2.64,
        sdduration: float = 0.11,
        xlim: Optional[Tuple[int, int]] = None,
    ) -> Figure:
        # TODO
        # need to simplify the difference of heatshock class and lethargus

        gridfig = plt.figure(figsize=(16, 6))
        if plotdata.lower() not in {"foq", "area"}:
            logging.info(f"{plotdata}: Not implemented")
            return gridfig

        sgs = {sg.groupname: sg for sg in self.samplegroups}
        label_axes = []
        for name, ind_num, idx in self.groups:
            sg = sgs[name]
            ax = gridfig.add_subplot(self.rownum, self.colnum, idx + 1)
            label_axes.append((idx, name, ax))
            if not sg.fullindlist:
                ax.annotate(
                    "No data",
                    xy=(0.5, 0.5),
                    xycoords="axes fraction",
                    fontsize=14,
                    horizontalalignment="center",
                    verticalalignment="center_baseline",
                    color="red",
                    alpha=0.8,
                    fontweight="bold",
                )
                continue

            ind = sg.fullindlist[ind_num - 1]
            if plotdata == "foq":
                ind.plot_foq(ax)
                if self.mode not in HEATSHOCK_LABELS and ind.has_valid_lethargus:
                    for lt in ind.letharguslist:
                        ind.plot_lethargus(ax, lt, sg.threshold)

            elif plotdata == "area":
                ax.plot(
                    ind.calc_raw_area(60)
                    .rolling(window=int(window / ind.interval), center=True)
                    .median(),
                    linestyle="-",
                    linewidth=0.5,
                    color="black",
                )

            if (
                self.mode not in HEATSHOCK_LABELS
                and ind.has_valid_lethargus
                and overlayparam is not None
            ):
                if overlayparam == "fq_duration":
                    lt = ind.letharguslist[0]
                    fontcolor = "black"
                    if lt.fq_duration > meanduration + 3 * sdduration:
                        fontcolor = "red"
                    elif lt.fq_duration < meanduration - 3 * sdduration:
                        fontcolor = "blue"
                    ax.annotate(
                        "{0:4.3}".format(lt.fq_duration),
                        xy=(lt.start, 0.1),
                        fontsize=10,
                        color=fontcolor,
                    )

            ax.set_ylim(0, 1.1)
            # ax.set_xlim(0, 5*60*60/interval)
            hrtick = np.arange(len(ind.foq) * ind.interval / 60 / 60).astype(int)
            ax.set_xticks(hrtick * 60 * 60 / ind.interval)
            if xlim is not None and len(xlim) == 2:
                ax.set_xlim(xlim)

            if self.mode in HEATSHOCK_LABELS:
                # TODO
                # this part is for manual imput the time window where to apply the heatshock
                # pre-capture for 1hr and 30 mins heatshock
                hs_start_hr = 1
                hs_end_hr = 1.5
                hs_start_frame = timedelta(hours=hs_start_hr).seconds / self.interval
                hs_end_frame = timedelta(hours=hs_end_hr).seconds / self.interval
                rect = patches.Rectangle(
                    (hs_start_frame, 0),
                    hs_end_frame - hs_start_frame,
                    1,
                    alpha=0.4,
                    color="red",
                    linewidth=0,
                )
                ax.set_xlim(0, timedelta(hours=6.5).seconds / self.interval)
                ax.add_patch(rect)
            ax.annotate(
                ind.label_str,
                xy=(0.01, 0.9),
                xycoords="axes fraction",
                fontsize=8,
                horizontalalignment="left",
                verticalalignment="bottom",
            )

        for ax in gridfig.get_axes():
            ax.spines["right"].set_visible(False)
            ax.spines["top"].set_visible(False)
            ax.set_xticklabels([])
            ax.set_yticklabels([])

        last_axes = [
            max(filter(lambda x: x[1] == name, label_axes), key=lambda x: x[0])
            for name in sgs.keys()
        ]

        # retrieve the tail of valid axes
        idx = max(
            (x[0] for x in last_axes if x[1] != UNDEFINE_LABEL),
            default=self.rownum * self.colnum - 1,
        )
        hrtick_label = list(map(str, hrtick))
        for i in range(len(hrtick_label)):
            if i % (1 + len(hrtick_label) // 10):
                hrtick_label[i] = ""

        ax = gridfig.get_axes()[idx]
        ax.set_xticklabels(hrtick_label)
        gridfig.tight_layout()
        gridfig.subplots_adjust(hspace=0.1)
        gridfig.subplots_adjust(wspace=0.01)
        return gridfig

    def to_pickle(self, compression: bool = True) -> "Project":
        suffix = ".pkl.gz" if compression else ".pkl"
        filename = self.datapath.with_suffix(suffix).name
        path = self.homepath.joinpath(filename)
        pd.to_pickle(self, path)
        logging.info(f"Save project at: {path}")
        return self

    @classmethod
    def from_pickle(cls, filepath: Union[str, Path]) -> "Project":
        if not Path(filepath).is_file():
            raise FileNotFoundError(filepath)
        return pd.read_pickle(filepath)

    @property
    def project_df(self) -> pd.DataFrame:
        if self._summary is None:
            self.make_project_df()
        return self._summary

    def make_project_df(self, *filters: Hashable) -> "Project":
        dfs: List[pd.DataFrame] = [
            sg.summarydf for sg in self.samplegroups if sg.summarydf is not None
        ]
        if not dfs:
            logging.warn(
                f"No summary_df found in {[sg.groupname for sg in self.samplegroups]}"
            )
            self._summary = pd.DataFrame()
            return self

        self._summary = pd.concat(
            dfs,
            ignore_index=True,
            axis=0,
        )
        if not filters:
            return

        col_filters = list(filters)
        if "groupname" in col_filters:
            col_filters.remove("groupname")
        # make sure groupname come first
        col_filters.insert(0, "groupname")

        remain = set(col_filters) - set(dfs[0].columns)
        if remain:
            logging.warn(f"Invalid columns name: {remain}")
            return self

        self._summary = pd.concat(
            dfs,
            ignore_index=True,
            axis=0,
        )[col_filters]
        return self

    def dotplots(
        self,
        params: str,
        *,
        uniquegroupnames: List[str] = None,
        labels: List[str] = None,
        ylim: Tuple[float, float] = None,
        figsize: Tuple[float, float] = None,
        thickness: float = 1,
        color: str = "black",
        outlier: str = "oc",
        size=3,
        coeff=0.3,
    ):
        if params not in self.project_df.columns:
            raise ValueError(
                f"{params} is not a valid lethargus parameters.\
                Only support {[n for n in self.project_df.columns]}"
            )
        groupname = self.project_df["groupname"].unique()
        uniquegroupnames = uniquegroupnames or [
            n for n in self.uniquegroupnames if n in groupname
        ]
        dfs = self.project_df.loc[
            self.project_df["groupname"].isin(uniquegroupnames), ["groupname", params]
        ]

        ymin = dfs[params].min()
        ymax = dfs[params].max()
        if ylim is not None:
            ymin, ymax = ylim
        # padding the margin with 5%
        y_margin = (ymax - ymin) * 0.05
        figsize = figsize or (0.65 * len(groupname), 4)
        fig = plt.figure(figsize=figsize)
        ax = fig.add_subplot(1, 1, 1)
        ax.spines["right"].set_visible(False)
        ax.spines["top"].set_visible(False)
        ax.yaxis.set_ticks_position("left")
        ax.xaxis.set_ticks_position("bottom")
        ax.set_ylim(ymin - y_margin, ymax + y_margin)
        ax.set_xlim(-0.5, len(uniquegroupnames) - 0.5)

        ncount = dfs.groupby("groupname")[params].count().to_dict()
        labelswithsize = [
            f"{name}\n{ncount.get(name, '')}" for name in uniquegroupnames
        ]
        if labels is not None and len(labels) >= dfs["groupname"].nunique():
            labelswithsize = [
                f"{name}\n{num}"
                for name, num in zip(labels, dfs.groupby("groupname").count())
            ]

        ax.set_xticks(np.arange(0, dfs["groupname"].nunique(), dtype=int))
        ax.set_xticklabels(labelswithsize)
        dfs.groupby("groupname")[params].apply(
            lambda df: add_jitter_plot(
                df.values,
                ax=ax,
                center_pos=uniquegroupnames.index(df.name),
                thickness=thickness,
                color=color,
                outlier=outlier,
                size=size,
                coeff=coeff,
                jitter_bin=5,
            )
        )
        fig.suptitle(params.capitalize())
        fig.tight_layout()
        return fig

    def create_summary_slide(self) -> "Project":
        slidename = self.homepath.joinpath(f"{self.datapath.stem}_summary.pptx")
        prs = Presentation()
        # landscape A4
        prs.slide_width = Cm(29.7)
        prs.slide_height = Cm(21.0)
        # empty layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        info = f"""Basic Parameters:
Threshold       : {self.foqthreshold}
number of column: {self.colnum}
number of row   : {self.rownum}
orientation     : {self.grouporder}
mode            : {self.mode}
"""

        shape = slide.shapes
        text_box = shape.add_textbox(
            left=Cm(-1.25), top=Cm(-0.5), height=Cm(0.5), width=Cm(7)
        )
        tf = text_box.text_frame

        p = tf.add_paragraph()
        p.text = info
        p.font.bold = True
        p.font.size = Pt(8)

        param = f"""File information:          
filename: {self.datapath.name}
date: {self.date}
experiment name: {self.expname}
number of experiment: {self.expnum}
frame interval (sec): {self.interval}
"""

        text_box = shape.add_textbox(
            left=Cm(4), top=Cm(-0.5), height=Cm(0.5), width=Cm(7)
        )
        tf = text_box.text_frame

        p = tf.add_paragraph()
        p.text = param
        p.font.bold = True
        p.font.size = Pt(8)

        # Add image
        top = 2.75
        for p in self.homepath.glob("*.png"):
            if "grid" in p.name.lower():
                shape.add_picture(str(p), Cm(0.25), Cm(top), width=Cm(14.2))
                top += 5.5
        shift = 0.25
        for dot in self.homepath.glob("*.png"):
            if "dot" in dot.name.lower():
                shape.add_picture(str(dot), Cm(shift), Cm(top), height=Cm(7))
                shift += 3.5

        shift = 0
        heatmaps = sorted(
            (
                f
                for f in self.homepath.glob("*.png")
                if "heatmap" in f.name.lower()
                and f.stem.split("_")[0] in self.uniquegroupnames
            ),
            key=lambda f: self.uniquegroupnames.index(f.stem.split("_")[0]),
        )
        if self.mode in HEATSHOCK_LABELS:
            heatmaps = [f for f in heatmaps if "_aligned" not in f.name.lower()]
        else:
            heatmaps = [f for f in heatmaps if "_aligned" in f.name.lower()]

        for f in heatmaps:
            shape.add_picture(str(f), Cm(14.75), Cm(2.75 + shift), height=Cm(3))
            shift += 3.2
        prs.save(slidename)
        return self

    def posthoc_screen_samplegroup(self, groupname: str) -> "Project":
        self._summary = None
        groupnames = {sg.groupname: sg for sg in self.samplegroups}
        if groupname not in groupnames:
            raise ValueError(f"Invalid groupname: {groupname}")
        self.homepath = self.homepath.parent.joinpath(
            datetime.now().strftime(r"%y%m%d_%H%M%S_manually_analyzed")
        )
        self.homepath.mkdir()
        sg = groupnames[groupname].set_directory(self.homepath)
        sg = (
            sg.manual_screen()
            .makedf(ltlist=sg.fullindlist)
            .savesummarydf()
            .makeltmatrix(align="head")
            .saveltmatrix(sg.ltfoqmatrix, "fqlt")
            .makeltmatrix(align="tail")
            .saveltmatrix(sg.ltfoqmatrixtail, "fqltaligntail")
            .saveafig(sg.makealignfigure(sg.ltfoqmatrix, True, "mean"), "fqaligned")
            .saveafig(
                sg.makealignfigure(sg.ltfoqmatrixtail, False, "mean"),
                "fqalignedtail",
            )
            .saveafig(sg.makeallfigureofareacropped(sg.indlist), "areacropped")
            .saveafig(sg.makeallfigureofarea(60, sg.fullindlist), "rawareafull")
            .saveafig(sg.plot_foq_heatmap(align=False), "foq_heatmap")
            .saveafig(sg.plot_foq_heatmap(align=True), "foq_heatmap_aligned")
        )
        return self
